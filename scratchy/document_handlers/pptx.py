from typing import Dict, Any
from .base import BaseDocumentHandler

class PPTXHandler(BaseDocumentHandler):
    """Handler for PowerPoint presentations (.pptx) using python-pptx."""

    def __init__(self, file_path: str):
        super().__init__(file_path)
        self._prs = None
        self._text = ""
        self._metadata = {}

    def load(self) -> None:
        """Load and parse the PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("python-pptx is not installed. Please install it via 'pip install python-pptx' to use PPTXHandler.")

        try:
            self._prs = Presentation(self.file_path)
            
            text_content = []
            for i, slide in enumerate(self._prs.slides):
                slide_text = []
                # Add slide header
                slide_text.append(f"--- Slide {i+1} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                text_content.append("\n".join(slide_text))
            
            self._text = "\n\n".join(text_content)
            
            # Extract core properties
            core_props = self._prs.core_properties
            self._metadata = {
                "author": core_props.author,
                "created": str(core_props.created),
                "modified": str(core_props.modified),
                "title": core_props.title,
                "subject": core_props.subject,
                "keywords": core_props.keywords,
                "last_modified_by": core_props.last_modified_by,
                "slide_count": len(self._prs.slides)
            }
            # Remove None values
            self._metadata = {k: v for k, v in self._metadata.items() if v is not None}

        except Exception as e:
            raise RuntimeError(f"Failed to load PPTX file {self.file_path}: {e}")

    def get_text(self) -> str:
        if self._prs is None:
            self.load()
        return self._text

    def get_metadata(self) -> Dict[str, Any]:
        if self._prs is None:
            self.load()
        return self._metadata
